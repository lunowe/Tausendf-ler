#!/usr/bin/env python3
"""Erzeugt docs/presentation/Tausendfuessler.pptx (python-pptx 1.0.x).

Alle Layouts werden aus Shapes/Textboxen auf der leeren Folienvorlage gebaut, damit das
Erscheinungsbild unabhaengig vom Standard-Theme ist. Farben und Motiv orientieren sich am
Corporate Design der Universitaet Leipzig (Basalt, Granat, Karneol, Aquamarin; Futura).

Aufruf:  python3 docs/presentation/build_pptx.py
Der Aufruf baut die Datei und prueft danach Folienanzahl, Notizen und (naeherungsweise) Textueberlauf.
"""
from __future__ import annotations

import math
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "Tausendfuessler.pptx"

# ---------------------------------------------------------------- Corporate Design
BASALT = RGBColor(0x26, 0x2A, 0x31)
GRANAT = RGBColor(0xB0, 0x2F, 0x2C)
KARNEOL = RGBColor(0xD8, 0x41, 0x3E)
AQUAMARIN = RGBColor(0x8A, 0xC2, 0xD1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6E, 0x73, 0x7A)
LIGHT = RGBColor(0xF2, 0xF3, 0xF4)
AQUA_LIGHT = RGBColor(0xE3, 0xF1, 0xF5)
RULE = RGBColor(0xD9, 0xDB, 0xDE)
FONT = "Futura"  # Fallback (Arial) regelt das Theme; Fonts koennen mit python-pptx nicht eingebettet werden.
MONO = "Menlo"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_W = SLIDE_W - 2 * MARGIN
CONTENT_TOP = Inches(1.55)
CONTENT_BOTTOM = Inches(6.75)

FOOTER_TEXT = "Tausendfüßler · EVA SS26 · Universität Leipzig"
EXPECTED_SLIDES = 15


# ---------------------------------------------------------------- Grundbausteine
def _apply_font(run, size, color=BASALT, bold=False, italic=False, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def _add_runs(paragraph, text_, size, color, bold, italic=False):
    """Mini-Markup: **fett** und `code`."""
    for part in re.split(r"(\*\*.+?\*\*|`.+?`)", text_):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            _apply_font(run, size, color, True, italic)
        elif part.startswith("`") and part.endswith("`"):
            run.text = part[1:-1]
            _apply_font(run, size - 1, color, bold, italic, name=MONO)
        else:
            run.text = part
            _apply_font(run, size, color, bold, italic)


def _set_bullet(paragraph, level, size):
    pPr = paragraph._p.get_or_add_pPr()
    indent = int(Pt(size * 1.1))
    pPr.set("marL", str(indent * (level + 1)))
    pPr.set("indent", str(-indent))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buClr", "a:buFont"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    bu_clr = pPr.makeelement(qn("a:buClr"), {})
    bu_clr.append(bu_clr.makeelement(qn("a:srgbClr"), {"val": "B02F2C" if level == 0 else "6E737A"}))
    pPr.append(bu_clr)
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "–" if level == 0 else "·"}))


def _no_bullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    if pPr.find(qn("a:buNone")) is None:
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def text(slide, x, y, w, h, items, size=18, color=BASALT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, bullets=False, spacing=1.12, space_after=6, italic=False,
         margin=Inches(0.05)):
    """Textbox ohne Auto-Fit. items: str oder Liste aus str | (str, level)."""
    if isinstance(items, str):
        items = [items]
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = margin
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(items):
        level = 0
        if isinstance(item, tuple):
            item, level = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        if bullets and item:
            _set_bullet(p, level, size)
        else:
            _no_bullet(p)
        _add_runs(p, item, size if level == 0 else size - 2, color, bold, italic)
    return shape


def box(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0, rotation=0):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.rotation = rotation
    s.text_frame.word_wrap = True
    return s


def panel(slide, x, y, w, h, fill=LIGHT):
    s = box(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    s.adjustments[0] = min(0.08, float(Inches(0.15)) / min(int(w), int(h)))
    return s


def label_box(slide, x, y, w, h, lines, fill, color=WHITE, size=14, line=None,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, bold_first=True, align=PP_ALIGN.CENTER):
    """Form mit zentriertem Text; erste Zeile fett, weitere Zeilen kleiner."""
    s = box(slide, x, y, w, h, fill=fill, line=line, shape=shape)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = min(0.15, float(Inches(0.12)) / min(int(w), int(h)))
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    if isinstance(lines, str):
        lines = [lines]
    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.05
        _no_bullet(p)
        _add_runs(p, line_text, size if i == 0 else size - 3, color, bold_first and i == 0)
    return s


def arrow(slide, x1, y1, x2, y2, color=BASALT, width=1.75, both=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if both:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def rule(slide, x1, y1, x2, y2, color=RULE, width=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def table(slide, x, y, w, col_widths, rows, size=12, row_h=Inches(0.34), first_col_bold=False):
    """Tabelle ohne Theme-Style: Kopfzeile Granat/weiss, Zeilen abwechselnd weiss/hellgrau, nur horizontale Linien."""
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, row_h * n_rows)
    tbl = shape.table
    tblPr = tbl._tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("firstRow", "0")
    style = tblPr.find(qn("a:tableStyleId"))
    if style is not None:
        tblPr.remove(style)
    total = sum(col_widths)
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = int(w * cw / total)
    for r, row in enumerate(rows):
        tbl.rows[r].height = row_h
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_header = r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRANAT if is_header else (LIGHT if r % 2 == 0 else WHITE)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            _no_bullet(p)
            _add_runs(p, str(value), size, WHITE if is_header else BASALT, is_header or (first_col_bold and c == 0))
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:lnL", "a:lnR"):
                ln = tcPr.makeelement(qn(tag), {"w": "0"})
                ln.append(ln.makeelement(qn("a:noFill"), {}))
                tcPr.append(ln)
            for tag in ("a:lnT", "a:lnB"):
                ln = tcPr.makeelement(qn(tag), {"w": "6350"})
                sf = ln.makeelement(qn("a:solidFill"), {})
                sf.append(sf.makeelement(qn("a:srgbClr"), {"val": "D9DBDE"}))
                ln.append(sf)
                tcPr.append(ln)
    return shape


def notes(slide, text_):
    slide.notes_slide.notes_text_frame.text = text_


def caption(slide, x, y, w, txt, size=11, align=PP_ALIGN.LEFT, color=GREY, h=Inches(0.3)):
    return text(slide, x, y, w, h, txt, size=size, color=color, align=align, space_after=0)


def bullets(slide, items, x=MARGIN, y=CONTENT_TOP, w=CONTENT_W, h=None, size=18, space_after=10):
    h = h if h is not None else CONTENT_BOTTOM - y
    return text(slide, x, y, w, h, items, size=size, bullets=True, space_after=space_after)


# ---------------------------------------------------------------- Folienrahmen
def _footer(slide, number):
    rule(slide, MARGIN, Inches(6.95), SLIDE_W - MARGIN, Inches(6.95))
    text(slide, MARGIN, Inches(7.0), Inches(8), Inches(0.35), FOOTER_TEXT, size=10, color=GREY, space_after=0)
    text(slide, SLIDE_W - MARGIN - Inches(1.5), Inches(7.0), Inches(1.5), Inches(0.35), str(number),
         size=10, color=GREY, align=PP_ALIGN.RIGHT, space_after=0)


def content_slide(prs, title, subtitle=None, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # leeres Layout
    # Zurueckhaltende Variante des CD-Motivs: schmaler Granat-Balken links oben mit Karneol-Fuss
    box(slide, 0, 0, Inches(0.14), Inches(1.35), fill=GRANAT)
    box(slide, 0, Inches(1.35), Inches(0.14), Inches(0.22), fill=KARNEOL)
    if section:
        text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.3), section.upper(), size=11, color=GRANAT,
             bold=True, space_after=0)
    text(slide, MARGIN, Inches(0.55), CONTENT_W, Inches(0.7), title, size=27, color=BASALT, bold=True,
         space_after=0, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        text(slide, MARGIN, Inches(1.2), CONTENT_W, Inches(0.35), subtitle, size=13, color=GREY, space_after=0)
    prs._counter = getattr(prs, "_counter", 0) + 1
    _footer(slide, prs._counter)
    return slide


# ---------------------------------------------------------------- Folien
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    prs._counter = getattr(prs, "_counter", 0) + 1
    # CD-Motiv: Granat-Grundform, von einem Karneol-Balken in 65 Grad gekreuzt
    box(slide, 0, 0, Inches(4.6), SLIDE_H, fill=GRANAT)
    box(slide, Inches(2.05), Inches(0.65), Inches(0.5), Inches(6.2), fill=KARNEOL, rotation=65)
    box(slide, Inches(4.6), Inches(5.9), Inches(0.3), Inches(1.6), fill=AQUAMARIN)
    text(slide, Inches(5.3), Inches(1.9), Inches(7.4), Inches(0.45), "PROTOTYP · ENTWICKLUNG VERTEILTER ANWENDUNGEN",
         size=12, color=GRANAT, bold=True, space_after=0)
    text(slide, Inches(5.3), Inches(2.35), Inches(7.4), Inches(1.3), "Tausendfüßler", size=54, color=BASALT,
         bold=True, space_after=0, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, Inches(5.3), Inches(3.7), Inches(7.4), Inches(0.9),
         "Ein verteilter Webcrawler, gesteuert über einen Telegram-Bot", size=22, color=BASALT, space_after=0)
    text(slide, Inches(5.3), Inches(5.2), Inches(7.4), Inches(1.2),
         ["Kanan Namazov · <TODO: Name Bot-Entwickler/in> · Luca Wegner",
          "Universität Leipzig · Institut für Wirtschaftsinformatik · Dr. Pascal Kovacs · SS 2026"],
         size=14, color=GREY, space_after=4)
    notes(slide, "Begrüßung und Vorstellung des Teams. Tausendfüßler ist unser Prototyp für das Modul "
                 "Entwicklung verteilter Anwendungen: ein Webcrawler, der auf mehrere Prozesse verteilt ist und "
                 "komplett über Telegram bedient wird. Der Vortrag folgt der vorgegebenen Struktur: Idee, Aufbau, "
                 "Entwicklungsprozess, Demo, Code-Details, danach Probleme und Verbesserungspotenzial.")


def slide_idea(prs):
    s = content_slide(prs, "Die Idee: Crawlen per Chat", section="1 · Idee",
                      subtitle="Ein Telegram-Chat als einzige Oberfläche für einen verteilten Crawler")
    bullets(s, [
        "Nutzer schickt **/crawl <URL> <Tiefe> [Filter …]** an den Bot",
        "Koordinator verteilt die URLs auf beliebig viele **Worker-Prozesse**",
        "Gefundene Seiten kommen als **Live-Stream** (URL, Titel, Textanfang) zurück in den Chat",
        "**/pause**, **/resume**, **/abort** greifen in laufende Aufträge ein – Ergebnisse bleiben erhalten",
        "Nach Abschluss: **Report** mit Seiten, Links, Fehlern und Dauer",
        "Dazu **/list**, **/status**, **/search** (Volltext über alle Ergebnisse) und **/stats**",
    ], w=Inches(7.1), size=17)
    cx, cw = Inches(8.2), Inches(4.4)
    panel(s, cx, CONTENT_TOP, cw, Inches(5.1))
    caption(s, cx + Inches(0.2), CONTENT_TOP + Inches(0.1), cw - Inches(0.4), "Telegram · @TausflerBot", size=10)

    def bubble(y, h, lines, mine):
        w = Inches(3.5)
        x = cx + cw - w - Inches(0.2) if mine else cx + Inches(0.2)
        label_box(s, x, y, w, h, lines, fill=AQUAMARIN if mine else WHITE, color=BASALT, size=11,
                  bold_first=False, align=PP_ALIGN.LEFT)

    bubble(CONTENT_TOP + Inches(0.45), Inches(0.42), ["/crawl https://example.org 2 blog"], True)
    bubble(CONTENT_TOP + Inches(1.0), Inches(0.6), ["Auftrag gestartet", "Job 3f9a… · Tiefe 2 · Filter: blog"], False)
    bubble(CONTENT_TOP + Inches(1.75), Inches(1.15),
           ["example.org/blog/ – Blog", "„Neuigkeiten aus dem Projekt …“",
            "example.org/blog/2026 – Archiv", "„Alle Beiträge des Jahres …“"], False)
    bubble(CONTENT_TOP + Inches(3.05), Inches(0.42), ["/status 3f9a"], True)
    bubble(CONTENT_TOP + Inches(3.6), Inches(1.15),
           ["Report – abgeschlossen", "437 Seiten · 5 210 Links · 3 Fehler", "Dauer 1 min 52 s"], False)
    notes(s, "Kurz die Nutzersicht zeigen: Der ganze Crawler wird mit acht Chat-Befehlen bedient. Rechts ein "
             "nachgestellter Chatverlauf: /crawl mit Start-URL, Tiefe und optionalem Filter, dann der Live-Stream "
             "der gefundenen Seiten und am Ende der Report. Die 437 Seiten stammen aus dem Livetest vom 29.08. "
             "gegen eine echte Sitemap; Links- und Fehlerzahl im Mockup sind Beispielwerte. Ursprung der Idee: "
             "docs/Skizze_v2.pdf.")


def slide_requirements(prs):
    s = content_slide(prs, "Die drei Modulanforderungen und wo sie erfüllt sind", section="1 · Idee")
    rows = [
        ["Anforderung", "Umsetzung im Prototyp", "Ort im Code"],
        ["Nebenläufig", "Thread-Pool je Worker (Größe = CPU-Kerne); ein Handler-Thread je Worker-Verbindung "
                        "im Koordinator; Bot pollt nebenläufig per @Scheduled",
         "worker/pool/CrawlExecutor · coordinator/socket/WorkerSocketServer · bot/service/ResultPoller"],
        ["Verteilt", "Koordinator – Worker über persistente TCP-Sockets (Port 9090, line-delimited JSON); "
                     "Worker sind eigenständige JVM-Prozesse, auch auf anderen Rechnern",
         "common/protocol/Message · coordinator/socket/* · worker/net/CoordinatorConnection"],
        ["Komponentenbasiert", "Koordinator als Spring-Boot-REST-Service (Port 8080, Spring MVC + Data JPA); "
                               "Bot als Spring-Boot-Client mit DI und Command-Registry",
         "coordinator/api/*Controller · bot/config/CommandRegistry · bot/service/CoordinatorClient"],
    ]
    table(s, MARGIN, CONTENT_TOP, CONTENT_W, [2.7, 5.7, 4.2], rows, size=13, row_h=Inches(0.5), first_col_bold=True)
    text(s, MARGIN, Inches(4.9), CONTENT_W, Inches(1.1),
         "Dazu aus der Skizze: Least-Work-First-Verteilung, thread-sichere URL-Deduplizierung, Absturz-Erkennung "
         "mit Re-Queue, Live-Ergebnisse in unter 2 s, Persistenz und Volltextsuche in Postgres, Retention-Cleanup "
         "beim Start (Standard 30 Tage).", size=14, color=GREY)
    notes(s, "Diese Folie beantwortet die Bestehensvoraussetzung direkt: Nebenläufigkeit, Verteilung und "
             "komponentenbasierte REST-Services. Wichtig ist der Hinweis, dass die Socket-Kommunikation zwischen "
             "Koordinator und Worker liegt und REST zwischen Bot und Koordinator – anders als in der ersten Skizze, "
             "in der Bot und Koordinator ebenfalls per Socket sprechen sollten. Die dritte Spalte nennt die "
             "Einstiegspunkte für die Code-Fragen.")


def slide_architecture(prs):
    s = content_slide(prs, "Architektur: vier Prozesstypen, zwei Schnittstellen", section="2 · Aufbau")
    y_mid, h = Inches(3.3), Inches(1.1)
    cy = y_mid + h // 2
    label_box(s, Inches(0.5), y_mid, Inches(1.7), h, ["Telegram", "externer Dienst"], fill=AQUAMARIN, color=BASALT)
    label_box(s, Inches(2.9), y_mid, Inches(2.1), h, ["Bot", "Spring Boot, kein Web-Server"], fill=BASALT)
    label_box(s, Inches(6.1), Inches(2.75), Inches(3.2), Inches(2.2),
              ["Koordinator", "Spring Boot: REST-API + TCP-Server", "Job-Queue · Frontier · Dedup",
               "Least-Work-First · Worker-Registry"], fill=GRANAT, size=16)
    worker_ys = [Inches(1.7), Inches(3.3), Inches(4.9)]
    for label, yy in zip(["Worker 1", "Worker 2", "Worker n"], worker_ys):
        label_box(s, Inches(10.7), yy, Inches(2.0), Inches(0.9), [label, "Plain Java · Thread-Pool"], fill=BASALT, size=14)
    label_box(s, Inches(6.7), Inches(5.5), Inches(1.8), Inches(1.05), ["Postgres", "Jobs · Seiten · Volltext"],
              fill=AQUAMARIN, color=BASALT, shape=MSO_SHAPE.CAN, size=13)

    arrow(s, Inches(2.2), cy, Inches(2.9), cy, both=True)
    caption(s, Inches(1.6), cy - Inches(0.75), Inches(1.9), "Bot-API, HTTPS Long Polling", size=9, align=PP_ALIGN.CENTER, h=Inches(0.45))
    arrow(s, Inches(5.0), cy, Inches(6.1), cy, both=True, color=GRANAT)
    caption(s, Inches(5.0), cy - Inches(0.5), Inches(1.1), "REST :8080", size=9, align=PP_ALIGN.CENTER, color=GRANAT)
    for yy in worker_ys:
        arrow(s, Inches(9.3), Inches(3.85), Inches(10.7), yy + Inches(0.45), both=True, color=GRANAT)
    caption(s, Inches(9.3), Inches(2.1), Inches(1.4), "TCP :9090", size=9, align=PP_ALIGN.CENTER, color=GRANAT)
    caption(s, Inches(9.2), Inches(2.35), Inches(1.6), "JSON je Zeile", size=9, align=PP_ALIGN.CENTER, color=GRANAT)
    arrow(s, Inches(7.6), Inches(4.95), Inches(7.6), Inches(5.5), both=True)
    caption(s, Inches(7.7), Inches(5.05), Inches(0.8), "JPA", size=9)

    text(s, MARGIN, Inches(5.7), Inches(5.6), Inches(1.05),
         ["Worker verbinden sich **zum** Koordinator und **holen** Arbeit ab (Pull).",
          "Der Koordinator **pusht** nur Steuersignale (PAUSE / RESUME / ABORT)."], size=13, color=GREY, space_after=2)
    text(s, Inches(9.5), Inches(6.0), Inches(3.3), Inches(0.7),
         "Modul `common`: `Message`-Records + Jackson als gemeinsames Protokoll", size=11, color=GREY, space_after=0)
    notes(s, "Vier Prozesstypen: Bot, Koordinator, n Worker und Postgres. Zwei eigene Schnittstellen: REST vom Bot "
             "zum Koordinator und TCP mit zeilenweisem JSON zwischen Koordinator und Workern; der Telegram-Verkehr "
             "ist ein externer Dienst. Wichtig für Nachfragen: Worker öffnen die Verbindung, deshalb reicht ein "
             "einziger offener Port am Koordinator, und Worker können auf beliebigen Rechnern laufen. "
             "Details in README.md und PROTOCOL.md.")


def slide_protocol(prs):
    s = content_slide(prs, "Zwei Protokolle: REST für den Bot, Socket für die Worker", section="2 · Aufbau")
    lw = Inches(5.95)
    caption(s, MARGIN, CONTENT_TOP - Inches(0.05), lw, "REST-API des Koordinators (Spring MVC, :8080)", size=12, color=GRANAT)
    rows = [
        ["HTTP", "Pfad", "Zweck"],
        ["POST", "/api/jobs", "Auftrag anlegen {url, maxDepth, filters, owner}"],
        ["GET", "/api/jobs?owner=", "eigene Aufträge (/list)"],
        ["GET", "/api/jobs/{id}", "Detail: Status, Tiefe, Seiten, Links, Fehler"],
        ["GET", "/api/jobs/{id}/results?afterSeq=", "Live-Stream-Seite, aufsteigend nach seq"],
        ["POST", "/api/jobs/{id}/pause | resume | abort", "Steuerung: 204, bei ungültigem Übergang 409"],
        ["GET", "/api/search?q= · /api/stats · /api/health", "Volltextsuche, Statistik, Startzeit"],
    ]
    table(s, MARGIN, CONTENT_TOP + Inches(0.3), lw, [1.0, 3.0, 3.4], rows, size=10, row_h=Inches(0.4))

    rx = MARGIN + lw + Inches(0.4)
    rw = CONTENT_W - lw - Inches(0.4)
    caption(s, rx, CONTENT_TOP - Inches(0.05), rw, "Socket-Protokoll Koordinator – Worker (TCP :9090)",
            size=12, color=GRANAT)
    rows = [
        ["Richtung", "type", "Bedeutung"],
        ["W → K", "REGISTER {workerId, threads}", "Handshake nach Connect"],
        ["W → K", "REQUEST_WORK {capacity}", "freie Slots melden (Pull)"],
        ["K → W", "WORK_PACKAGE {jobId, depth, urls[]}", "≤ min(capacity, 2 · threads) URLs"],
        ["K → W", "NO_WORK {retryAfterMs}", "später erneut fragen"],
        ["W → K", "PAGE_RESULT {url, title, links[], error}", "ein Ergebnis je URL"],
        ["K → W", "JOB_SIGNAL {PAUSE | RESUME | ABORT}", "Push, jederzeit"],
    ]
    table(s, rx, CONTENT_TOP + Inches(0.3), rw, [0.9, 3.0, 2.4], rows, size=10, row_h=Inches(0.4))
    text(s, MARGIN, Inches(5.2), CONTENT_W, Inches(1.2), [
        "Beide Schnittstellen sind JSON; die Socket-Typen sind Records hinter einem `sealed interface` in `common` "
        "(Jackson, Diskriminator `type`).",
        "Der Live-Stream ist ein Cursor: `seq` steigt pro Job streng monoton, der Bot fragt jede Sekunde "
        "`results?afterSeq=<zuletzt gesehen>` ab.",
    ], size=13, color=GREY, bullets=True, space_after=4)
    notes(s, "Links die REST-Schnittstelle, rechts die Socket-Nachrichtentypen. Auf Nachfrage den Ablauf aus "
             "PROTOCOL.md erklären: REGISTER, dann REQUEST_WORK, WORK_PACKAGE, PAGE_RESULT im Kreis; PAUSE und "
             "ABORT werden vom Koordinator gepusht, ohne dass der Worker fragt. Code: "
             "common/.../protocol/Message.java (sealed interface mit Records) und ProtocolJson.java.")


def slide_coordinator(prs):
    s = content_slide(prs, "Koordinator: Zustand pro Job, ein Thread pro Worker", section="5 · Code-Details",
                      subtitle="coordinator/…/crawl/JobRuntime · crawl/Scheduler · socket/WorkerConnectionHandler · service/ResultService")
    lw = Inches(6.4)
    bullets(s, [
        "**JobRuntime** (In-Memory je laufendem Job): Frontier `ConcurrentSkipListMap<Tiefe, Queue>`, "
        "Dedup `ConcurrentHashMap.newKeySet()`, In-Flight-Map URL → (Worker, Tiefe)",
        "Verbundoperationen `takeWork`, `complete`, `requeue`, `isFinished` sind `synchronized`: ein Job gilt nie "
        "als fertig, während eine URL zwischen Frontier und In-Flight ist",
        "**Scheduler**: Worker pullen, sobald Slots frei sind – der am wenigsten belastete Worker fragt am "
        "häufigsten (Least-Work-First); Jobs werden Round-Robin bedient",
        "**WorkerSocketServer**: Accept-Thread + `CachedThreadPool`, ein **WorkerConnectionHandler** je "
        "Verbindung; Read-Timeout 60 s ⇒ Absturz erkannt ⇒ `requeue(workerId)`",
        "**ResultService.handle()**: unter dem Runtime-Monitor `seq++`, Seite speichern, Links in die Frontier, "
        "Zähler, COMPLETED-Prüfung",
    ], w=lw, size=14, space_after=8)
    rx = MARGIN + lw + Inches(0.4)
    rw = CONTENT_W - lw - Inches(0.4)
    steps = [
        ("PAGE_RESULT", "vom Handler-Thread des Workers"),
        ("synchronized (runtime)", "Monitor des JobRuntime"),
        ("complete(url) · nextSeq()", "In-Flight raus, Cursor hoch"),
        ("pages.save(PageEntity)", "JPA → Postgres"),
        ("offerLinks(depth, links)", "visited.add() – first one wins"),
        ("isFinished() → COMPLETED", "Frontier leer und nichts in flight"),
    ]
    yy = CONTENT_TOP
    for i, (head, sub) in enumerate(steps):
        hot = i in (1, 4)
        label_box(s, rx, yy, rw, Inches(0.6), [head, sub], fill=GRANAT if hot else LIGHT,
                  color=WHITE if hot else BASALT, size=13, align=PP_ALIGN.LEFT)
        if i < len(steps) - 1:
            arrow(s, rx + rw // 2, yy + Inches(0.6), rx + rw // 2, yy + Inches(0.78), color=GREY, width=1.25)
        yy += Inches(0.78)
    caption(s, rx, yy, rw, "Pfad eines Ergebnisses in ResultService.handle()", size=10)
    notes(s, "Zum Vorführen öffnen: coordinator/src/main/java/.../crawl/JobRuntime.java (Frontier, visited, "
             "inFlight, synchronized takeWork/complete/requeue/isFinished), crawl/Scheduler.java (assign), "
             "socket/WorkerConnectionHandler.java (Zeilenschleife, Read-Timeout, requeue bei EOF) und "
             "service/ResultService.java (handle). Kernaussage: Dedup ist atomar über visited.add(), die "
             "Verbundoperationen sind über den Monitor des JobRuntime serialisiert. Tests: JobRuntimeTest, "
             "SchedulerTest, WorkerSocketServerTest, EndToEndTest.")


def slide_worker(prs):
    s = content_slide(prs, "Worker: Plain Java, Thread-Pool, Pause-Warteschleife", section="5 · Code-Details",
                      subtitle="worker/…/WorkerClient · CrawlTask · pool/CrawlExecutor · crawler/PageFetcher · crawler/HtmlExtractor")
    lw = Inches(6.4)
    bullets(s, [
        "**Kein Spring**: `java -jar worker.jar --coordinator host:9090`, Start in unter 1 s",
        "**CrawlExecutor**: `Executors.newFixedThreadPool(Kerne)`; `submit()` liefert "
        "`CompletableFuture<CrawlOutcome>` (`CrawlSuccess` | `CrawlFailure`, sealed)",
        "**WorkerClient**: Main-Schleife sendet `REQUEST_WORK` mit `threads − inFlight` freien Slots, höchstens "
        "eine offene Anfrage; eigener Reader-Thread verarbeitet die Antworten",
        "**CrawlTask.awaitPermission()**: prüft vor jedem Abruf `jobControl` (ConcurrentHashMap Job → Signal); "
        "bei PAUSE Warteschleife (200 ms), bei ABORT `CancellationException`",
        "**PageFetcher** (java.net.http, 10 s Timeout, folgt Redirects) + **HtmlExtractor** (jsoup: Titel, Text, "
        "Links, Sitemap-`<loc>`); relative Links werden gegen die **finale** URL aufgelöst",
        "**Reconnect** mit Backoff 1 s → 10 s; laufende Tasks werden bei Verbindungsverlust verworfen",
    ], w=lw, size=14, space_after=8)
    rx = MARGIN + lw + Inches(0.4)
    rw = CONTENT_W - lw - Inches(0.4)
    iw = rw - Inches(0.6)
    panel(s, rx, CONTENT_TOP, rw, Inches(5.0))
    caption(s, rx + Inches(0.3), CONTENT_TOP + Inches(0.12), iw, "Threads in einem Worker-Prozess", size=11, color=GRANAT)
    label_box(s, rx + Inches(0.3), CONTENT_TOP + Inches(0.5), iw, Inches(0.65),
              ["main: runSession()", "REQUEST_WORK, wait/notify-Wakeup"], fill=BASALT, size=13)
    label_box(s, rx + Inches(0.3), CONTENT_TOP + Inches(1.3), iw, Inches(0.65),
              ["coordinator-reader", "readLoop → handle(): WORK_PACKAGE, JOB_SIGNAL"], fill=BASALT, size=13)
    for i in range(3):
        label_box(s, rx + Inches(0.3) + i * (iw // 3), CONTENT_TOP + Inches(2.25), iw // 3 - Inches(0.08), Inches(0.9),
                  [f"pool-{i + 1}", "fetch → extract"], fill=GRANAT, size=12)
    caption(s, rx + Inches(0.3), CONTENT_TOP + Inches(3.25), iw,
            "ein Crawl-Thread je CPU-Kern (availableProcessors)", size=10)
    label_box(s, rx + Inches(0.3), CONTENT_TOP + Inches(3.7), iw, Inches(0.95),
              ["whenComplete → PAGE_RESULT", "inFlight − 1, Wakeup der Main-Schleife"], fill=WHITE, color=BASALT,
              size=13, line=GREY)
    notes(s, "Öffnen: worker/.../WorkerClient.java (runSession, readMessages, submit/complete), CrawlTask.java "
             "(awaitPermission mit Pause-Schleife) und pool/CrawlExecutor.java. Der Crawl-Kern (PageFetcher, "
             "HtmlExtractor, CrawlExecutor) stammt von Kanan; die Socket-Schleife und die Signalbehandlung kamen "
             "bei der Integration dazu. Für die NFA 'Thread-Pool liefert alles': "
             "CrawlExecutorTest.handlesHighVolumeConcurrentCrawls – 100 URLs auf 4 Threads, exakt 100 Ergebnisse.")


def slide_bot(prs):
    s = content_slide(prs, "Bot: Long Polling, Command-Registry, Ergebnis-Poller", section="5 · Code-Details",
                      subtitle="bot/…/TausflerBot · config/CommandRegistry · commands/*CommandHandler · service/ResultPoller · CoordinatorClient")
    lw = Inches(6.4)
    bullets(s, [
        "**TausflerBot** erbt `TelegramLongPollingBot`: Updates kommen per Long Polling, kein eingehender Port nötig",
        "**CommandRegistry**: alle `CommandHandler`-Beans werden per DI injiziert; Klassenname "
        "`CrawlCommandHandler` → Befehl `/crawl` (Konvention statt Konfiguration)",
        "**CoordinatorClient** kapselt die REST-Aufrufe (`RestClient`) und übersetzt 400/404/409 in lesbare Chat-Meldungen",
        "**ResultPoller** (`@Scheduled` alle 1 s): Map Job → (chatId, `lastSeq`); holt `results?afterSeq=` seitenweise "
        "und bündelt mehrere Seiten je Telegram-Nachricht (Rate-Limit ≈ 1 Nachricht/s pro Chat)",
        "Status **vor** den Ergebnissen abfragen: ist der Job terminal, sind alle Seiten bis dahin garantiert "
        "ausgeliefert – dann Report senden und abbestellen",
    ], w=lw, size=14, space_after=8)
    rx = MARGIN + lw + Inches(0.4)
    rw = CONTENT_W - lw - Inches(0.4)
    rows = [
        ["Befehl", "Handler → REST"],
        ["/crawl url tiefe [filter]", "POST /api/jobs, dann subscribe()"],
        ["/list", "GET /api/jobs?owner=chatId"],
        ["/status id", "GET /api/jobs/{id}"],
        ["/pause · /resume · /abort id", "POST /api/jobs/{id}/…"],
        ["/search text", "GET /api/search?q="],
        ["/stats", "GET /api/stats"],
        ["/help", "CommandRegistry.getAllCommands()"],
    ]
    table(s, rx, CONTENT_TOP, rw, [2.5, 3.0], rows, size=11, row_h=Inches(0.4))
    caption(s, rx, CONTENT_TOP + Inches(3.4), rw, "Job-IDs dürfen als Präfix angegeben werden (z. B. die ersten 4 Zeichen).", size=10, h=Inches(0.5))
    notes(s, "Öffnen: bot/.../TausflerBot.java (onUpdateReceived, Dispatch über die Registry), "
             "config/CommandRegistry.java (Namenskonvention), commands/CrawlCommandHandler.java und "
             "service/ResultPoller.java (pollJob mit seq-Cursor und Report). Das Bot-Skelett mit Registry und "
             "Handlern hat unser Bot-Teammitglied gebaut; der seq-Cursor, das Bündeln der Nachrichten und die "
             "Fehlerübersetzung kamen beim Integrieren dazu. Tests: ResultPollerTest, CoordinatorClientTest.")


def slide_concurrency(prs):
    s = content_slide(prs, "Nebenläufigkeit und Verteilung auf einen Blick", section="5 · Code-Details")
    rows = [
        ["Mechanismus", "Primitive / Technik", "Klasse"],
        ["Paralleles Crawlen", "FixedThreadPool (n = Kerne), CompletableFuture-Verkettung", "worker/pool/CrawlExecutor, CrawlTask"],
        ["Ein Thread je Worker-Verbindung", "CachedThreadPool + Accept-Thread, Read-Timeout 60 s", "coordinator/socket/WorkerSocketServer"],
        ["URL-Dedup, atomar", "ConcurrentHashMap.newKeySet().add() – „first one wins“", "coordinator/crawl/JobRuntime.offerLinks"],
        ["Frontier und In-Flight konsistent", "synchronized-Methoden auf dem JobRuntime-Monitor", "JobRuntime.takeWork / complete / requeue"],
        ["Status-Übergang ohne Lost Update", "synchronized (runtime) um load-update-save", "service/JobService.underRuntimeLock, ResultService.handle"],
        ["Pause / Abort im Worker", "ConcurrentHashMap<Job, Signal>, volatile Flags, wait/notify", "worker/CrawlTask.awaitPermission, WorkerClient"],
        ["Live-Stream", "AtomicLong seq je Job, @Scheduled-Poller mit Cursor", "JobRuntime.nextSeq, bot/service/ResultPoller"],
        ["Absturz-Recovery", "EOF / Timeout im Handler ⇒ requeue aller In-Flight-URLs des Workers", "WorkerConnectionHandler, JobRuntime.requeue"],
        ["Verteilung", "TCP-Sockets (line-delimited JSON) · REST (Spring MVC) · JPA / Postgres", "common/protocol, coordinator/api, persistence"],
    ]
    table(s, MARGIN, CONTENT_TOP, CONTENT_W, [3.0, 5.0, 4.4], rows, size=12, row_h=Inches(0.46), first_col_bold=True)
    notes(s, "Das ist die Spickzettel-Folie für die Detailfragen: jede Zeile nennt das Problem, das verwendete "
             "Java-Primitive und die Klasse. Bei Nachfragen zur Thread-Sicherheit auf JobRuntime verweisen: die "
             "Collections sind concurrent, die Verbundoperationen synchronized. Beim Lost Update erklären, dass "
             "Statuswechsel und Ergebnisverarbeitung denselben Monitor nehmen, weil beide die Job-Zeile per "
             "load-update-save schreiben.")


def slide_process(prs):
    s = content_slide(prs, "Entwicklungsprozess und Aufgabenteilung", section="3 · Entwicklungsprozess",
                      subtitle="Zeitachse aus dem Git-Log (2026): zuerst drei getrennte Repositories, am Ende ein Multimodul-Build")
    ty = Inches(2.7)
    x0, x1 = MARGIN + Inches(0.4), SLIDE_W - MARGIN - Inches(0.4)
    rule(s, x0, ty, x1, ty, color=GREY, width=1.5)
    milestones = [
        ("07.05.", "Thema und erster Plan", "Luca", 0.00),
        ("28.05.", "Worker-Crawl-Kern: Fetch, jsoup, Thread-Pool, Tests", "Kanan", 0.2),
        ("05.07.", "Skizze v2 · Lasttest im Worker", "Kanan, Team", 0.4),
        ("08.–27.08.", "Bot: Handler, Registry, Poller, DTOs", "TODO: Name (Git „debian“)", 0.63),
        ("29.08.", "Konsolidierung, Koordinator, Protokoll, Integration, Lasttest, Livetest", "Luca", 0.9),
    ]
    for i, (date, what, who, frac) in enumerate(milestones):
        x = int(x0 + (x1 - x0) * frac)
        box(s, x - Inches(0.09), ty - Inches(0.09), Inches(0.18), Inches(0.18), fill=GRANAT, shape=MSO_SHAPE.OVAL)
        w = Inches(2.5)
        bx = min(max(x - w // 2, MARGIN), SLIDE_W - MARGIN - w)
        if i % 2 == 0:
            text(s, bx, ty - Inches(1.1), w, Inches(0.95), [f"**{date}** {what}", who], size=11,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, space_after=0)
        else:
            text(s, bx, ty + Inches(0.15), w, Inches(0.95), [f"**{date}** {what}", who], size=11,
                 align=PP_ALIGN.CENTER, space_after=0)
    cy = Inches(4.0)
    ch = Inches(2.7)
    cw = (CONTENT_W - Inches(0.6)) // 3
    cards = [
        ("Kanan Namazov", "Worker-Crawl-Kern",
         ["PageFetcher (java.net.http), HtmlExtractor (jsoup)", "CrawlExecutor mit Thread-Pool und Futures",
          "Unit- und Integrationstests inkl. Lasttest (100 URLs / 4 Threads)"]),
        ("<TODO: Name eintragen>", "Telegram-Bot (Git-Autor „debian“)",
         ["TausflerBot, CommandRegistry, alle CommandHandler", "DTOs, CoordinatorClient, MessageSender",
          "ResultPoller mit Live-Stream (erste Fassung)"]),
        ("Luca Wegner", "Koordinator, Protokoll, Integration",
         ["JobRuntime, Scheduler, TCP-Server, REST-API, JPA, Suche", "Worker-Socket-Schleife, Pause/Abort, Reconnect",
          "End-to-End-Test, Lasttest-Modul, NFA-Report, Fixes – mit KI-Tooling (Claude Code, Codex)"]),
    ]
    for i, (name, role, items) in enumerate(cards):
        x = MARGIN + i * (cw + Inches(0.3))
        panel(s, x, cy, cw, ch)
        box(s, x, cy, Inches(0.09), ch, fill=GRANAT if i == 2 else AQUAMARIN)
        text(s, x + Inches(0.2), cy + Inches(0.1), cw - Inches(0.3), Inches(0.4), name, size=15, bold=True, space_after=0)
        text(s, x + Inches(0.2), cy + Inches(0.48), cw - Inches(0.3), Inches(0.35), role, size=12, color=GRANAT, space_after=0)
        text(s, x + Inches(0.2), cy + Inches(0.88), cw - Inches(0.3), Inches(1.75), items, size=11, bullets=True, space_after=3)
    notes(s, "Ehrlich darstellen, wer was gemacht hat: Kanan hat im Mai/Juli den Crawl-Kern des Workers mit Tests "
             "gebaut, das Bot-Teammitglied im August das Bot-Skelett mit Registry, Handlern und erstem Poller, "
             "Luca am 29.08. den Koordinator, das gemeinsame Protokoll, die Integration in ein Multimodul-Repo, "
             "den Lasttest und die Fixes – dabei mit KI-Werkzeugen (Claude Code, Codex) als Pair-Programmer. "
             "Vor dem Vortrag den Platzhalter TODO durch den Namen ersetzen (git log --format='%an %s').")


def slide_nfa(prs):
    s = content_slide(prs, "NFA-Ergebnisse: gemessen, nicht geschätzt", section="2 · Aufbau",
                      subtitle="loadtest-Modul (Plain Java) gegen Koordinator + Postgres 16, Lauf vom 29.08.2026 · docs/NFA-Report.md")
    rows = [
        ["NFA (aus der Skizze)", "Grenzwert", "Messwert", "Szenario"],
        ["Koordinator startbereit", "< 15 s", "3,4 s (JVM-Start → ApplicationReady)", "/api/health startupSeconds"],
        ["Statusabfrage bei 20 Aufträgen", "< 0,2 s", "p50 7 ms · p95 13 ms · max 98 ms", "300 Abfragen, 0 Überschreitungen"],
        ["Anteil ohne internen Fehler", "> 99,9 %", "100 % (0 × 5xx bei 1 200 Anfragen)", "60 s Mix gültig/ungültig, 654 × 4xx erwartet"],
        ["Durchsatz 2 vs. 1 Worker", "≥ +60 %", "+92 % (20,5 → 39,2 Seiten/s)", "600 Seiten, 100 ms Abrufzeit, 4 Threads je Worker"],
        ["Live-Ergebnis nach Abruf", "< 2 s", "p50 254 ms · p95 507 ms · max 517 ms", "125 Seiten, Poll-Intervall 500 ms"],
        ["URL-Dedup atomar", "0 Duplikate", "342 Seiten, 0 doppelt bei 18 076 gemeldeten Links", "500 Seiten, 50 gemeinsame Hubs, 2 Worker"],
        ["Thread-Pool vollständig", "100 / 100", "exakt 100 CrawlOutcomes", "CrawlExecutorTest (JUnit)"],
        ["Mehrkernauslastung", "> 1 Thread", "3 Handler-Threads (1 130 / 537 / 309 Ergebnisse)", "Thread-IDs in den Logs"],
    ]
    table(s, MARGIN, CONTENT_TOP, CONTENT_W, [3.0, 1.3, 4.2, 3.9], rows, size=12, row_h=Inches(0.42), first_col_bold=True)
    text(s, MARGIN, Inches(5.95), CONTENT_W, Inches(0.75),
         "Ehrlicher Befund: ohne Netzlatenz (0 ms Abrufzeit) nur +28 % (87 → 112 Seiten/s), weil der Koordinator "
         "die Ergebnisse je Job serialisiert und pro Seite synchron in die DB schreibt – siehe Verbesserungspotenzial.",
         size=12, color=GREY)
    notes(s, "Alle Zahlen stammen aus docs/NFA-Report.md, erzeugt vom loadtest-Modul gegen den echten Stack mit "
             "Postgres. Das Skalierungsergebnis ehrlich einordnen: mit realistischer Abruflatenz skaliert es fast "
             "linear, ohne Latenz bremst der serialisierte Ergebnispfad im Koordinator. Auf Nachfrage: "
             "loadtest/.../scenario/ThroughputScenario.java und DedupScenario.java.")


def slide_demo(prs):
    s = content_slide(prs, "Demo-Ablauf", section="4 · Demo",
                      subtitle="Live über Telegram; Backup: Screenshots und Logs aus dem Livetest vom 29.08.2026 (logs/)")
    lw = Inches(6.6)
    steps = [
        ("1", "Stack starten", "docker compose up · coordinator.jar · 2 × worker.jar · bot.jar – Log zeigt Start in ≈ 3 s und zwei REGISTER"),
        ("2", "/crawl <Sitemap-URL> 1", "Live-Stream läuft im Chat; im Koordinator-Log gehen Pakete an beide Worker"),
        ("3", "/pause · /status · /resume", "Worker stoppen nach den laufenden Abrufen; Status PAUSED; danach geht es weiter"),
        ("4", "Worker-Prozess beenden (kill)", "Handler erkennt EOF, URLs werden neu verteilt, Job läuft mit einem Worker zu Ende"),
        ("5", "/abort auf einem zweiten Job", "Ergebnisse bleiben erhalten – /search und /stats zeigen sie"),
    ]
    yy = CONTENT_TOP
    for num, head, sub in steps:
        label_box(s, MARGIN, yy, Inches(0.55), Inches(0.75), num, fill=GRANAT, size=18)
        text(s, MARGIN + Inches(0.7), yy - Inches(0.02), lw - Inches(0.7), Inches(0.85), [f"**{head}**", sub],
             size=13, space_after=0)
        yy += Inches(0.98)
    rx = MARGIN + lw + Inches(0.5)
    rw = CONTENT_W - lw - Inches(0.5)
    panel(s, rx, CONTENT_TOP, rw, Inches(4.9), fill=AQUA_LIGHT)
    text(s, rx + Inches(0.25), CONTENT_TOP + Inches(0.15), rw - Inches(0.5), Inches(4.6), [
        "**Fallback, wenn Netz oder Telegram streiken**",
        "Screenshots des Livetests (437 Seiten, alle Befehle) – TODO: in docs/presentation/ ablegen",
        "`logs/coordinator.out`, `logs/worker1.out`, `logs/bot.out` aus demselben Lauf",
        "`mvn test` mit `EndToEndTest`: REST → Koordinator → 2 echte Worker → lokale Test-Site → H2",
        "**Vorbereitung**: `.env` mit Bot-Token, Postgres auf Port 5433, Telegram-Chat auf dem Beamer",
    ], size=13, bullets=True, space_after=8)
    notes(s, "Reihenfolge der Demo: Stack starten, Crawl mit Live-Stream, Pause/Resume, dann einen Worker hart "
             "beenden, um die Recovery zu zeigen, zuletzt Abort mit erhaltenen Ergebnissen. Wenn etwas hängt, "
             "sofort auf Screenshots und Logs ausweichen; der EndToEndTest läuft ohne Netz und ohne Docker gegen H2. "
             "Vorher prüfen: Bot-Token gültig, Port 5433 frei, keine fremden Worker verbunden.")


def slide_problems(prs):
    s = content_slide(prs, "Probleme und Hürden", section="Rückblick")
    rows = [
        ["Hürde", "Was passiert ist", "Lösung"],
        ["Zwei konkurrierende Entwürfe", "Worker- und Bot-Repo entstanden unabhängig, mit eigenen DTOs und Annahmen "
                                        "(Bot erwartete REST, Skizze sagte Socket)",
         "Konsolidierung in ein Multimodul-Repo, PROTOCOL.md als Vertrag, Modul common"],
        ["Lost Update beim Job-Status", "/pause und ein gleichzeitiges PAGE_RESULT schrieben dieselbe Zeile per "
                                       "load-update-save – der neue Status ging verloren",
         "Beide Pfade unter dem Monitor des JobRuntime (JobService.underRuntimeLock)"],
        ["Worker forderte zu viel an", "Jede fertige Seite löste ein neues REQUEST_WORK aus – ein Worker bekam "
                                      "mehr URLs als Threads",
         "Höchstens eine offene Anfrage; Kapazität = Threads − inFlight"],
        ["Redirects", "Relative Links wurden gegen die angeforderte statt der finalen URL aufgelöst – falsche Ziele",
         "FetchResult.finalUrl aus HttpResponse.uri()"],
        ["Telegram-Rate-Limit", "Eine Nachricht je Seite lief in „429 Too Many Requests“",
         "Seiten je Poll-Runde bündeln, Nachrichtenlänge begrenzen"],
        ["Infrastruktur", "Lokales Homebrew-Postgres blockierte Port 5432; ddl-auto konnte eine neue NOT-NULL-Spalte "
                          "nicht anlegen",
         "Compose auf Port 5433; Spaltendefault in JobEntity"],
    ]
    table(s, MARGIN, CONTENT_TOP, CONTENT_W, [2.4, 5.6, 4.2], rows, size=11, row_h=Inches(0.5), first_col_bold=True)
    notes(s, "Die größte Hürde war organisatorisch: drei Teilprojekte mit unterschiedlichen Annahmen, die erst am "
             "Ende zusammengeführt wurden. Technisch am lehrreichsten war der Lost Update beim Job-Status, ein "
             "klassisches Nebenläufigkeitsproblem zwischen REST-Thread und Worker-Handler-Thread. Die Zeilen "
             "entsprechen den Commits 'Review-Fixes: Status-Race, Handler-Robustheit, Live-Stream, Redirect-Basis-URL' "
             "und 'End-to-End-Test und Fix für Worker-Überzuteilung'.")


def slide_improvements(prs):
    s = content_slide(prs, "Verbesserungspotenzial", section="Ausblick")
    lw = (CONTENT_W - Inches(0.5)) // 2
    text(s, MARGIN, CONTENT_TOP, lw, Inches(0.35), "Leistung und Skalierung", size=14, color=GRANAT, bold=True, space_after=0)
    bullets(s, [
        "Ergebnispfad im Koordinator entkoppeln: DB-Schreiben asynchron oder gebündelt statt synchron unter dem "
        "Job-Monitor (heute +28 % statt +92 % ohne Netzlatenz)",
        "Server-Sent Events oder WebSocket statt 1-s-Polling im Bot",
        "Expliziter Heartbeat statt 60 s Read-Timeout; robots.txt und Rate-Limit je Domain",
        "Frontier bei sehr großen Jobs in die DB auslagern (heute In-Memory)",
    ], x=MARGIN, y=CONTENT_TOP + Inches(0.4), w=lw, size=14)
    rx = MARGIN + lw + Inches(0.5)
    text(s, rx, CONTENT_TOP, lw, Inches(0.35), "Funktion und Struktur", size=14, color=GRANAT, bold=True, space_after=0)
    bullets(s, [
        "/search nur über eigene Jobs (Owner-Filter) und Ranking statt LIKE-/tsvector-Mix",
        "currentDepth als Frontier-Tiefe statt „höchste beantwortete Tiefe“",
        "JSON-Export je Job (aus der Skizze übernommen, noch offen)",
        "Bot-Handler mit gemeinsamer Basisklasse für das Argument-Parsing; Koordinator-Konfiguration "
        "über @ConfigurationProperties",
        "Security bewusst ausgeklammert – ein Token je Chat wäre der nächste Schritt",
    ], x=rx, y=CONTENT_TOP + Inches(0.4), w=lw, size=14)
    notes(s, "Das wichtigste Verbesserungspotenzial ist messbar: die Serialisierung des Ergebnispfads im Koordinator "
             "begrenzt den Durchsatz ohne Netzlatenz. Danach die Live-Ausspielung per SSE statt Polling und ein "
             "echter Heartbeat. Auf der Funktionsseite sind Owner-Filter für /search, die Semantik von currentDepth "
             "und der JSON-Export aus der Skizze noch offen.")


def slide_end(prs):
    s = content_slide(prs, "Zusammenfassung und Fragen", section="Abschluss")
    tiles = [
        ("4", "Prozesstypen", "Bot · Koordinator · Worker · Postgres"),
        ("2", "Schnittstellen", "REST :8080 · TCP :9090"),
        ("+92 %", "Durchsatz", "2 Worker gegenüber 1 Worker"),
        ("0", "Duplikate", "bei 18 076 gemeldeten Links"),
    ]
    tw = (CONTENT_W - 3 * Inches(0.3)) // 4
    for i, (num, head, sub) in enumerate(tiles):
        x = MARGIN + i * (tw + Inches(0.3))
        panel(s, x, CONTENT_TOP, tw, Inches(2.0))
        box(s, x + Inches(0.1), CONTENT_TOP, tw - Inches(0.2), Inches(0.07), fill=KARNEOL)
        text(s, x, CONTENT_TOP + Inches(0.25), tw, Inches(0.9), num, size=40, color=GRANAT, bold=True,
             align=PP_ALIGN.CENTER, space_after=0)
        text(s, x, CONTENT_TOP + Inches(1.15), tw, Inches(0.4), head, size=15, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        text(s, x, CONTENT_TOP + Inches(1.5), tw, Inches(0.5), sub, size=10, color=GREY, align=PP_ALIGN.CENTER, space_after=0)
    bullets(s, [
        "Nebenläufig: Thread-Pools in Worker und Koordinator, atomare Dedup, Monitor je Job",
        "Verteilt: persistente TCP-Sockets mit zeilenweisem JSON, Worker auf beliebigen Rechnern",
        "Komponentenbasiert: Spring-Boot-REST-Koordinator, Spring-Boot-Bot mit DI, Plain-Java-Worker",
        "Alle acht NFAs aus der Skizze gemessen und erfüllt; Livetest über Telegram mit 437 Seiten",
    ], y=Inches(3.95), h=Inches(2.0), size=15, space_after=6)
    text(s, MARGIN, Inches(6.1), CONTENT_W, Inches(0.5), "Fragen? – Code: README.md, PROTOCOL.md, docs/NFA-Report.md",
         size=16, color=GRANAT, bold=True, space_after=0)
    notes(s, "Abschluss in einem Satz: ein verteilter Crawler mit klarer Rollenteilung der Prozesse, dessen "
             "nicht-funktionale Anforderungen nicht behauptet, sondern gemessen sind. Dann zur Fragerunde "
             "überleiten und die Code-Detail-Folien 6 bis 9 als Einstieg für die Nachfragen anbieten.")


# ---------------------------------------------------------------- Aufbau + Pruefung
def _set_theme_fonts(prs):
    """Theme-Schriften auf Futura setzen, damit auch Reste des Standard-Themes konsistent aussehen."""
    theme_part = prs.slide_master.part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
    xml = theme_part.blob.decode("utf-8")
    xml = re.sub(r'<a:latin typeface="[^"]*"', f'<a:latin typeface="{FONT}"', xml)
    theme_part._blob = xml.encode("utf-8")


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    _set_theme_fonts(prs)
    for builder in (slide_title, slide_idea, slide_requirements, slide_architecture, slide_protocol,
                    slide_coordinator, slide_worker, slide_bot, slide_concurrency, slide_process, slide_nfa,
                    slide_demo, slide_problems, slide_improvements, slide_end):
        builder(prs)
    return prs


# Naeherung fuer die Ueberlaufpruefung: mittlere Zeichenbreite ~0.55 em (Futura ist breit), Zeilenhoehe 1.2 em
CHAR_W = 0.55
LINE_H = 1.2


def _frame_height_pt(tf, width_emu, default_size):
    total = 0.0
    width_pt = Emu(width_emu).pt
    for p in tf.paragraphs:
        size = default_size
        for r in p.runs:
            if r.font.size is not None:
                size = r.font.size.pt
        txt = "".join(r.text for r in p.runs)
        pPr = p._p.pPr
        indent = Emu(int(pPr.get("marL"))).pt if pPr is not None and pPr.get("marL") else 0
        usable = max(width_pt - indent, 1)
        chars_per_line = max(int(usable / (size * CHAR_W)), 1)
        lines = max(1, math.ceil(len(txt) / chars_per_line)) if txt else 1
        spacing = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
        total += lines * size * LINE_H * spacing
        if p.space_after is not None:
            total += p.space_after.pt
    return total


def check(path: Path, expected_slides: int) -> list[str]:
    prs = Presentation(str(path))
    problems = []
    if len(prs.slides) != expected_slides:
        problems.append(f"Folienanzahl {len(prs.slides)} statt {expected_slides}")
    for idx, slide in enumerate(prs.slides, start=1):
        if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            problems.append(f"Folie {idx}: keine Sprechernotizen")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                tf = shape.text_frame
                inner_w = shape.width - tf.margin_left - tf.margin_right
                inner_h = Emu(shape.height - tf.margin_top - tf.margin_bottom).pt
                need = _frame_height_pt(tf, inner_w, 18)
                if need > inner_h * 1.05 + 2:
                    problems.append(f"Folie {idx}: Text „{tf.text[:40]}…“ braucht ≈{need:.0f} pt, Box hat {inner_h:.0f} pt")
            if shape.has_table:
                tbl = shape.table
                total = 0.0
                for row in tbl.rows:
                    row_need = Emu(row.height).pt
                    for c, cell in enumerate(row.cells):
                        inner_w = tbl.columns[c].width - cell.margin_left - cell.margin_right
                        need = _frame_height_pt(cell.text_frame, inner_w, 12) + Emu(cell.margin_top + cell.margin_bottom).pt
                        row_need = max(row_need, need)
                    total += row_need
                allowed = Emu(CONTENT_BOTTOM - shape.top).pt
                if total > allowed:
                    problems.append(f"Folie {idx}: Tabelle braucht ≈{total:.0f} pt, Platz bis Inhaltsende {allowed:.0f} pt")
            if shape.rotation == 0 and (shape.left + shape.width > SLIDE_W + Inches(0.01)
                                        or shape.top + shape.height > SLIDE_H + Inches(0.01)):
                problems.append(f"Folie {idx}: Form „{shape.name}“ ragt über den Folienrand")
    return problems


def main():
    prs = build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    problems = check(OUT, expected_slides=EXPECTED_SLIDES)
    print(f"geschrieben: {OUT} ({len(prs.slides)} Folien)")
    if problems:
        print("Pruefung meldet:")
        for p in problems:
            print("  -", p)
        sys.exit(1)
    print("Pruefung ok: Folienanzahl, Notizen, kein erkennbarer Textueberlauf, nichts ragt über den Rand.")


if __name__ == "__main__":
    main()
